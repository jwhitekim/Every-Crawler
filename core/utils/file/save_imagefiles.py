import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt  

def load_image_files(image_folder):
    """ 폴더 내 모든 이미지 파일을 screenshot_XXX 순서대로 정렬하여 반환 """
    image_files = [
        os.path.join(image_folder, file)
        for file in os.listdir(image_folder)
        if file.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif'))
    ]

    # 숫자 기준 정렬 (screenshot_XXX.png → XXX 기준)
    image_files.sort(key=lambda x: int(os.path.basename(x).split('_')[1].split('.')[0]))

    return image_files


def PPTsave(image_folder, output_file="output.pptx"):
    # PowerPoint 생성
    presentation = Presentation()

    image_files = load_image_files(image_folder)

    presentation.slide_width = Inches(13.33)  # 가로 크기
    presentation.slide_height = Inches(7.5)   # 세로 크기

    # 이미지 파일을 슬라이드에 추가
    for image in image_files:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # 빈 슬라이드
        
        # 슬라이드 크기 가져오기
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        
        # 이미지 추가 (슬라이드 전체를 채우도록)
        slide.shapes.add_picture(image, 0, 0, width=slide_width, height=slide_height)

    # 결과 저장
    presentation.save(output_file)
    print(f"[INFO] PPT file has been saved: {output_file}")

def PDFsave(image_folder, output_pdf="output.pdf"):
    image_files = load_image_files(image_folder)

    # 이미지 로드 및 변환
    images = []
    for file in image_files:
        img = Image.open(file)
        if img.mode != "RGB":  # 이미지가 RGB 모드가 아닐 경우 변환
            img = img.convert("RGB")
        images.append(img)

    # PDF 생성
    if images:
        images[0].save(output_pdf, save_all=True, append_images=images[1:])
        print(f"[INFO] PDF creation complete : {output_pdf}")
    else:
        print("[DEBUG] There are no images to convert to PDF.")